import io
import json
import os
from typing import Generator

import fal_client
from subconscious import Subconscious

from .config import settings
from .models import (
    CharacterMap,
    ComicPanel,
    ComicResult,
    ComicScript,
    ConceptExtraction,
    ImagePrompts,
    ParsedLecture,
    PipelineEvent,
    SlideContent,
)

STEP_NAMES = [
    "Parsing slideshow",
    "Extracting concepts",
    "Mapping characters",
    "Writing panel scripts",
    "Generating image prompts",
    "Drawing panels",
]


def _emit(
    step: int, status: str, data: dict | None = None, error: str | None = None
) -> PipelineEvent:
    return PipelineEvent(
        step=step,
        step_name=STEP_NAMES[step - 1],
        status=status,
        data=data,
        error=error,
    )


# ── Step 1: Parse uploaded file ──


def parse_pptx(file_bytes: bytes) -> ParsedLecture:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slides.append(SlideContent(slide_number=i, text="\n".join(texts)))
    title = slides[0].text.split("\n")[0] if slides else ""
    full_text = "\n\n".join(f"[Slide {s.slide_number}]\n{s.text}" for s in slides)
    return ParsedLecture(title=title, slides=slides, full_text=full_text)


def parse_pdf(file_bytes: bytes) -> ParsedLecture:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=file_bytes, filetype="pdf")
    slides = []
    for i, page in enumerate(doc, 1):
        text = page.get_text()
        slides.append(SlideContent(slide_number=i, text=text))
    doc.close()
    title = slides[0].text.split("\n")[0] if slides else ""
    full_text = "\n\n".join(f"[Page {s.slide_number}]\n{s.text}" for s in slides)
    return ParsedLecture(title=title, slides=slides, full_text=full_text)


def parse_file(file_bytes: bytes, filename: str) -> ParsedLecture:
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pptx":
        return parse_pptx(file_bytes)
    elif ext == "pdf":
        return parse_pdf(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")


# ── Steps 2-5: Subconscious calls ──


def _sub_run(client: Subconscious, instructions: str, answer_format: type):
    run = client.run(
        engine=settings.engine,
        input={
            "instructions": instructions,
            "answerFormat": answer_format,
        },
        options={"await_completion": True},
    )
    if not run.result or not run.result.answer:
        raise RuntimeError("No result from Subconscious")
    return answer_format.model_validate_json(run.result.answer)


def extract_concepts(client: Subconscious, lecture: ParsedLecture) -> ConceptExtraction:
    instructions = f"""You are a concept extraction expert. Analyze these lecture notes and extract
the key concepts (3-8), identify the central tension or conflict that could drive a comic narrative,
and outline a narrative arc for a 5-panel comic strip.

Lecture content:
{lecture.full_text}"""
    return _sub_run(client, instructions, ConceptExtraction)


def map_characters(
    client: Subconscious, concepts: ConceptExtraction
) -> CharacterMap:
    instructions = f"""You are a creative character designer for educational comics.
Turn these abstract concepts into vivid, memorable characters for a comic strip.
Each character should visually embody the concept they represent.
Also choose a consistent art style for the entire comic.

Concepts: {concepts.model_dump_json()}
Central tension: {concepts.central_tension}
Topic: {concepts.lecture_topic}"""
    return _sub_run(client, instructions, CharacterMap)


def write_script(
    client: Subconscious,
    concepts: ConceptExtraction,
    characters: CharacterMap,
) -> ComicScript:
    instructions = f"""You are a comic book script writer specializing in educational content.
Write a 5-panel comic script using these characters and concepts.
Use the central tension to create a compelling mini-story with a beginning, middle, and end.
Each panel needs a setting, which characters appear, dialogue, action description, and mood.

Characters: {characters.model_dump_json()}
Narrative arc: {concepts.narrative_arc}
Central tension: {concepts.central_tension}
Topic: {concepts.lecture_topic}"""
    return _sub_run(client, instructions, ComicScript)


def generate_prompts(
    client: Subconscious,
    script: ComicScript,
    characters: CharacterMap,
) -> ImagePrompts:
    char_visuals = json.dumps(
        [{"name": c.name, "visual": c.visual_description} for c in characters.characters]
    )
    instructions = f"""You are an expert at writing image generation prompts for AI image models.
For each panel in this comic script, write a detailed image generation prompt.
Include: composition, character appearances (reference the visual descriptions), lighting, mood,
and the art style: {characters.art_style}

Each prompt MUST start with the art style for consistency across panels.
Avoid requesting text, lettering, or speech bubbles in the image.
Also write a negative_prompt for things to avoid.

Comic script: {script.model_dump_json()}
Character visual references: {char_visuals}"""
    return _sub_run(client, instructions, ImagePrompts)


# ── Step 6: fal.ai image generation ──


def generate_image(prompt: str, negative_prompt: str = "") -> str:
    # Ensure FAL_KEY is set in the environment for fal_client
    if settings.fal_key:
        os.environ["FAL_KEY"] = settings.fal_key

    full_negative = "text, words, letters, speech bubbles, watermark"
    if negative_prompt:
        full_negative = f"{negative_prompt}, {full_negative}"

    response = fal_client.run(
        "fal-ai/flux/dev",
        arguments={
            "prompt": prompt,
            "image_size": "landscape_16_9",
            "num_images": 1,
            "enable_safety_checker": True,
            "num_inference_steps": 28,
            "guidance_scale": 3.5,
        },
    )
    return response["images"][0]["url"]


# ── Full pipeline generator ──


def run_pipeline(file_bytes: bytes, filename: str) -> Generator[PipelineEvent, None, None]:
    client = Subconscious(api_key=settings.subconscious_api_key)

    # Step 1: Parse
    yield _emit(1, "started")
    try:
        lecture = parse_file(file_bytes, filename)
        yield _emit(1, "completed", data={"slide_count": len(lecture.slides), "title": lecture.title})
    except Exception as e:
        yield _emit(1, "error", error=str(e))
        return

    # Step 2: Concept Extraction
    yield _emit(2, "started")
    try:
        concepts = extract_concepts(client, lecture)
        yield _emit(2, "completed", data=concepts.model_dump())
    except Exception as e:
        yield _emit(2, "error", error=str(e))
        return

    # Step 3: Character Mapping
    yield _emit(3, "started")
    try:
        characters = map_characters(client, concepts)
        yield _emit(3, "completed", data=characters.model_dump())
    except Exception as e:
        yield _emit(3, "error", error=str(e))
        return

    # Step 4: Script Writing
    yield _emit(4, "started")
    try:
        script = write_script(client, concepts, characters)
        yield _emit(4, "completed", data=script.model_dump())
    except Exception as e:
        yield _emit(4, "error", error=str(e))
        return

    # Step 5: Image Prompt Generation
    yield _emit(5, "started")
    try:
        prompts = generate_prompts(client, script, characters)
        yield _emit(5, "completed", data=prompts.model_dump())
    except Exception as e:
        yield _emit(5, "error", error=str(e))
        return

    # Step 6: Image Generation (fal.ai)
    yield _emit(6, "started")
    try:
        panels: list[ComicPanel] = []
        for i, (panel_script, img_prompt) in enumerate(
            zip(script.panels, prompts.prompts)
        ):
            image_url = generate_image(
                prompt=img_prompt.prompt,
                negative_prompt=img_prompt.negative_prompt,
            )
            panel = ComicPanel(
                panel_number=panel_script.panel_number,
                setting=panel_script.setting,
                characters=panel_script.characters_present,
                dialogue=panel_script.dialogue,
                action=panel_script.action,
                image_url=image_url,
                image_prompt=img_prompt.prompt,
            )
            panels.append(panel)
            # Yield progress so frontend can show images as they arrive
            yield _emit(
                6,
                "progress",
                data={
                    "panel_number": panel.panel_number,
                    "image_url": image_url,
                    "panels_done": i + 1,
                    "panels_total": 5,
                },
            )

        result = ComicResult(title=script.title, panels=panels)
        yield _emit(6, "completed", data=result.model_dump())
    except Exception as e:
        yield _emit(6, "error", error=str(e))
        return
